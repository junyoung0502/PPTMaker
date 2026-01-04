import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import base64

# Streamlit 앱 설정
st.set_page_config(
    page_title="Euodia lyrics PPT",
    page_icon="📊",
    layout="wide"
)

# -------------------------------------------------------------------------
# CSS 스타일: 플로팅 다운로드 버튼
# -------------------------------------------------------------------------
st.markdown("""
<style>
    /* 다운로드 버튼 우측 하단 고정 */
    [data-testid="stDownloadButton"] {
        position: fixed;
        bottom: 30px;
        right: 30px;
        z-index: 9999;
    }
    
    /* 버튼 스타일 */
    [data-testid="stDownloadButton"] button {
        background-color: #FF4B4B;
        color: white;
        width: auto;
        padding: 15px 30px;
        border-radius: 50px;
        border: none;
        box-shadow: 0 4px 14px rgba(0,0,0,0.3);
        font-size: 1.2rem;
        font-weight: bold;
        transition: transform 0.1s;
    }
    
    [data-testid="stDownloadButton"] button:hover {
        background-color: #FF2B2B;
        transform: scale(1.05);
        color: white;
    }
    
    [data-testid="stDownloadButton"] button:active {
        transform: scale(0.95);
    }
</style>
""", unsafe_allow_html=True)

# -------------------------------------------------------------------------
# 함수 정의
# -------------------------------------------------------------------------

@st.cache_data(show_spinner=False)
def create_ppt_binary(title, text_content, font_name, font_size, x_pos, y_pos, text_alignment):
    """PPT 생성 및 바이너리 반환 (캐싱 적용)"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    align_map = {"왼쪽": PP_ALIGN.LEFT, "가운데": PP_ALIGN.CENTER, "오른쪽": PP_ALIGN.RIGHT}
    ppt_align = align_map.get(text_alignment, PP_ALIGN.CENTER)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]
    
    # 제목 슬라이드
    title_slide = prs.slides.add_slide(blank_slide_layout)
    fill = title_slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    title_box = title_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(2.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_run = title_paragraph.runs[0]
    title_run.font.name = font_name
    title_run.font.size = Pt(font_size + 6)
    title_run.font.color.rgb = RGBColor(255, 255, 255)
    title_run.font.bold = True
    
    # 내용 슬라이드
    slides_content = parse_text_to_slides(text_content)
    FIXED_WIDTH, FIXED_HEIGHT = 11.33, 3.0
    
    for slide_text in slides_content:
        if slide_text.strip():
            slide = prs.slides.add_slide(blank_slide_layout)
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)
            
            textbox = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(FIXED_WIDTH), Inches(FIXED_HEIGHT))
            text_frame = textbox.text_frame
            text_frame.text = slide_text
            text_frame.word_wrap = True
            
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = ppt_align
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.bold = True
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def parse_text_to_slides(text_content):
    lines = [line.strip() for line in text_content.split('\n') if line.strip()]
    slides, current_slide = [], []
    for line in lines:
        if line == '---':
            if current_slide: slides.append('\n'.join(current_slide)); current_slide = []
        else:
            current_slide.append(line)
            if len(current_slide) == 2: slides.append('\n'.join(current_slide)); current_slide = []
    if current_slide: slides.append('\n'.join(current_slide))
    return slides

def split_text_preview(text_content):
    return parse_text_to_slides(text_content)

# -------------------------------------------------------------------------
# UI 구성
# -------------------------------------------------------------------------

st.title("Euodia lyrics PPT")
st.markdown("---")
st.markdown("""
### 사용법
1. **PPT 제목**과 **내용**을 입력하세요.
2. **설정**을 조절하고 **미리보기**로 확인하세요.
3. 우측 하단의 **'📥 PPT 다운로드'** 버튼을 누르면 즉시 파일이 저장됩니다.
""")
st.markdown("---")

col1, col2, col3 = st.columns([1, 0.8, 1])

with col1:
    st.subheader("📝 입력")
    ppt_title = st.text_input("PPT 제목", value="꽃들도")
    default_text = """이 곳에 생명샘 솟아나
눈물 골짝 지나갈 때에
머잖아 열매 맺히고
웃음 소리 넘쳐나리라
이 곳에 생명샘 솟아나
---
눈물 골짝 지나갈 때에
머잖아 열매 맺히고
웃음 소리 넘쳐나리라
꽃들도 구름도
바람도 넓은 바다도
찬양하라 찬양하라 예수를
하늘을 울리며 노래해
나의 영혼아
은혜의 주 은혜의 주 은혜의 주
---
그날에 하늘이 열리고
모든 이가 보게 되리라
마침내 꽃들이 피고
영광의 주가 오시리라"""
    text_content = st.text_area("텍스트 내용", value=default_text, height=400)

with col2:
    st.subheader("⚙️ 설정")
    font_options = {'맑은 고딕': 'Malgun Gothic', '굴림': 'Gulim', '돋움': 'Dotum', '바탕': 'Batang', 'Arial': 'Arial'}
    selected_font_display = st.selectbox("폰트 선택", list(font_options.keys()))
    selected_font = font_options[selected_font_display]
    
    st.markdown("**텍스트 정렬**")
    alignment_option = st.radio("정렬", ["왼쪽", "가운데", "오른쪽"], index=1, horizontal=True, label_visibility="collapsed")
    
    st.markdown("**글자 크기 (pt)**")
    font_size = st.number_input("글자 크기", 10, 200, 54, 2)
    
    st.markdown("---")
    st.markdown("**텍스트 위치 (인치)**")
    col_x, col_y = st.columns(2)
    with col_x: x_pos = st.number_input("가로 (X)", 0.0, 13.0, 1.0, 0.1)
    with col_y: y_pos = st.number_input("세로 (Y)", 0.0, 7.0, 0.5, 0.1)
    
    st.info(f"설정: {font_size}pt / ({x_pos}, {y_pos})")

with col3:
    st.subheader("👀 미리보기")
    
    if text_content.strip():
        slides = split_text_preview(text_content)
        total_slides = len(slides)
        st.info(f"총 {total_slides + 1}개의 슬라이드")
        
        # 더보기 상태 관리 (Session State)
        if 'show_all_preview' not in st.session_state:
            st.session_state.show_all_preview = False

        # 제목 슬라이드
        st.markdown(f"<div style='background:black;color:white;padding:20px;text-align:center;border-radius:10px;margin-bottom:20px;font-family:{selected_font};'><h3 style='margin:0;'>{ppt_title}</h3></div>", unsafe_allow_html=True)
        
        # 비율 계산
        left_p, top_p = (x_pos/13.33)*100, (y_pos/7.5)*100
        w_p, h_p = (11.33/13.33)*100, (3.0/7.5)*100
        css_align = {"왼쪽":"left", "가운데":"center", "오른쪽":"right"}[alignment_option]
        
        # 보여줄 슬라이드 개수 결정
        display_slides = slides if st.session_state.show_all_preview else slides[:3]
        
        # 슬라이드 렌더링
        for idx, slide_content in enumerate(display_slides):
            formatted = slide_content.replace('\n', '<br>')
            p_size = min(font_size * 0.4, 20)
            st.markdown(f"""
            <div style='position: relative; background: black; width: 100%; padding-top: 56.25%; border-radius: 10px; margin-bottom: 15px; border: 1px solid #333;'>
                <div style='position: absolute; left: 10px; top: 10px; color: #666; font-size: 12px;'>Slide {idx+2}</div>
                <div style='position: absolute; left: {left_p}%; top: {top_p}%; width: {w_p}%; height: {h_p}%; border: 2px dashed #ffff00; display: flex;'>
                    <div style='font-family: "{selected_font}"; font-size: {p_size}px; text-align: {css_align}; width: 100%; color: white; font-weight: bold;'>
                        {formatted}
                    </div>
                </div>
            </div>
            """, unsafe_allow_html=True)

        # 더보기/접기 버튼 로직
        if total_slides > 3:
            if st.session_state.show_all_preview:
                if st.button("🔼 접기 (3개만 보기)", use_container_width=True):
                    st.session_state.show_all_preview = False
                    st.rerun()
            else:
                remaining = total_slides - 3
                if st.button(f"🔽 더보기 ({remaining}개 더 있음)", use_container_width=True):
                    st.session_state.show_all_preview = True
                    st.rerun()

# -------------------------------------------------------------------------
# 플로팅 다운로드
# -------------------------------------------------------------------------
if text_content.strip():
    try:
        ppt_binary = create_ppt_binary(ppt_title, text_content, selected_font, font_size, x_pos, y_pos, alignment_option)
        st.download_button(
            label="📥 PPT 다운로드",
            data=ppt_binary,
            file_name=f"{ppt_title}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception as e:
        st.error(f"오류: {e}")

# 사이드바
with st.sidebar:
    st.header("도움말")
    st.markdown("---")
    st.markdown("**문의**: jylee0005@gmail.com")